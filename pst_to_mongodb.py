# pst_to_mongodb.py
# Imports all folders from a .pst / .ost / Google Takeout mbox-as-pst file into MongoDB.
#
# REQUIREMENTS (install once):
#   pip install libpff-python pymongo python-dateutil
#
# USAGE:
#   python pst_to_mongodb.py --pst "C:\path\to\file.pst" --mongo "mongodb://localhost:27017" --db mydb
#
# OPTIONAL FLAGS:
#   --collection   MongoDB collection name (default: pst_items)
#   --batch        Insert batch size (default: 100)
#   --verbose      Print each item as it is inserted

import argparse
import sys
import os
import re

# Ensure Unicode characters print safely on Windows consoles
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
import hashlib
import datetime
import email.parser
import email.policy
import traceback

try:
    import pypff
except ImportError:
    print("ERROR: libpff-python is not installed.")
    print("Run:  pip install libpff-python")
    sys.exit(1)

try:
    import mimetypes
    from pymongo import MongoClient, UpdateOne
    from pymongo.errors import BulkWriteError
    import gridfs as gridfs_module
except ImportError:
    print("ERROR: pymongo is not installed.")
    print("Run:  pip install pymongo")
    sys.exit(1)

try:
    from dateutil import parser as dateparser
    HAS_DATEUTIL = True
except ImportError:
    HAS_DATEUTIL = False

# Optional text-extraction libraries (graceful fallback if missing)
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx as _docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl as _openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation as _Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False


# ---------------------------------------------------------------------------
# Attachment text extraction
# ---------------------------------------------------------------------------
_TEXT_LIMIT = 25_000   # chars per attachment stored in MongoDB

def extract_text_from_bytes(data: bytes, filename: str, content_type: str = "") -> str:
    """
    Extract plain text from attachment binary data.
    Returns empty string if the type is unsupported or extraction fails.
    """
    if not data:
        return ""
    import io
    ext = os.path.splitext((filename or "").lower())[1]
    ct  = (content_type or "").lower()

    # ---- Plain text / CSV / log ----
    if ext in (".txt", ".csv", ".tsv", ".log", ".md", ".ini", ".xml", ".json") \
            or ct.startswith("text/plain"):
        try:
            return data.decode("utf-8", errors="replace")[:_TEXT_LIMIT]
        except Exception:
            return ""

    # ---- HTML ----
    if ext in (".html", ".htm") or "html" in ct:
        try:
            if HAS_BS4:
                return BeautifulSoup(data, "html.parser").get_text(" ", strip=True)[:_TEXT_LIMIT]
            # Fallback: stdlib
            from html.parser import HTMLParser
            class _P(HTMLParser):
                def __init__(self): super().__init__(); self.parts = []
                def handle_data(self, d): self.parts.append(d)
            p = _P(); p.feed(data.decode("utf-8", errors="replace"))
            return " ".join(p.parts)[:_TEXT_LIMIT]
        except Exception:
            return ""

    # ---- PDF ----
    if ext == ".pdf" or "pdf" in ct:
        if not HAS_PDF:
            return ""
        try:
            return (pdf_extract_text(io.BytesIO(data)) or "")[:_TEXT_LIMIT]
        except Exception:
            return ""

    # ---- DOCX ----
    if ext == ".docx" or "wordprocessingml" in ct:
        if not HAS_DOCX:
            return ""
        try:
            doc = _docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text)[:_TEXT_LIMIT]
        except Exception:
            return ""

    # ---- XLSX ----
    if ext in (".xlsx", ".xls") or "spreadsheet" in ct or "excel" in ct:
        if not HAS_XLSX:
            return ""
        try:
            wb = _openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.extend(str(v) for v in row if v is not None)
            return " ".join(parts)[:_TEXT_LIMIT]
        except Exception:
            return ""

    # ---- PPTX ----
    if ext == ".pptx" or "presentationml" in ct:
        if not HAS_PPTX:
            return ""
        try:
            prs = _Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)[:_TEXT_LIMIT]
        except Exception:
            return ""

    # ---- RTF (basic strip) ----
    if ext == ".rtf" or "rtf" in ct:
        try:
            text = data.decode("latin-1", errors="replace")
            # Strip RTF control words
            text = re.sub(r"\\[a-z]+\-?\d*[ ]?", " ", text)
            text = re.sub(r"[{}\\]", "", text)
            return text[:_TEXT_LIMIT]
        except Exception:
            return ""

    return ""


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def safe_str(value, fallback=""):
    if value is None:
        return fallback
    if isinstance(value, bytes):
        return value.decode("utf-8", errors="replace")
    return str(value)


def parse_date(value):
    if value is None:
        return None
    if isinstance(value, datetime.datetime):
        return value
    s = str(value).strip()
    if not s:
        return None
    if HAS_DATEUTIL:
        try:
            return dateparser.parse(s)
        except Exception:
            pass
    return None


def make_id(folder_path, index):
    raw = f"{folder_path}||{index}"
    return hashlib.sha256(raw.encode("utf-8", errors="replace")).hexdigest()


def parse_transport_headers(message):
    """
    Parse the RFC-2822 transport headers string that pypff exposes.
    Returns a dict with keys: from_addr, to_addrs, cc_addrs, date, message_id.
    """
    result = {
        "from_addr": "",
        "to_addrs": [],
        "cc_addrs": [],
        "date": None,
        "message_id": "",
    }
    try:
        raw = message.transport_headers
        if not raw:
            return result
        raw = safe_str(raw)

        # Use stdlib email parser -- it handles folded headers, encoding, etc.
        p = email.parser.HeaderParser(policy=email.policy.compat32)
        msg = p.parsestr(raw)

        result["from_addr"] = safe_str(msg.get("From", "")).strip()
        result["message_id"] = safe_str(msg.get("Message-ID", "")).strip()

        to_raw = msg.get("To", "")
        if to_raw:
            result["to_addrs"] = [a.strip() for a in to_raw.split(",") if a.strip()]

        cc_raw = msg.get("Cc", "")
        if cc_raw:
            result["cc_addrs"] = [a.strip() for a in cc_raw.split(",") if a.strip()]

        date_raw = msg.get("Date", "")
        if date_raw and HAS_DATEUTIL:
            try:
                result["date"] = dateparser.parse(date_raw)
            except Exception:
                pass
    except Exception:
        pass
    return result


# MAPI property tag constants for attachment metadata
_PR_DISPLAY_NAME       = 0x3001
_PR_ATTACH_EXTENSION   = 0x3703
_PR_ATTACH_FILENAME    = 0x3704
_PR_ATTACH_LONG_FILENAME = 0x3707
_PR_ATTACH_MIME_TAG    = 0x370E


def _mapi_string(att, *tags):
    """Read a UTF-8 string from a pypff attachment's MAPI record sets by tag ID(s)."""
    try:
        for rs in att.record_sets:
            for entry in rs.entries:
                if entry.entry_type in tags:
                    for getter in ("get_data_as_string", "get_data_as_utf8_string"):
                        try:
                            v = getattr(entry, getter)()
                            if v and str(v).strip():
                                return str(v).strip()
                        except Exception:
                            pass
    except Exception:
        pass
    return ""


def _att_filename(att, index):
    """
    Extract attachment filename from MAPI record_set properties.
    Priority: long filename → short filename → display name → mime-based fallback.
    """
    for tags in (
        (_PR_ATTACH_LONG_FILENAME,),   # PR_ATTACH_LONG_FILENAME  e.g. "report.pdf"
        (_PR_ATTACH_FILENAME,),         # PR_ATTACH_FILENAME        8.3 name
        (_PR_DISPLAY_NAME,),            # PR_DISPLAY_NAME
    ):
        name = _mapi_string(att, *tags)
        if name:
            return name

    # Last resort: build from extension or MIME type
    ext  = _mapi_string(att, _PR_ATTACH_EXTENSION)   # e.g. ".pdf"
    mime = _mapi_string(att, _PR_ATTACH_MIME_TAG)     # e.g. "image/png"
    if not ext and mime:
        import mimetypes as _mt
        ext = _mt.guess_extension(mime) or ""
    return f"attachment_{index}{ext}"


def _att_mime(att, filename):
    """Get MIME type: first try MAPI PR_ATTACH_MIME_TAG, then guess from filename."""
    mime = _mapi_string(att, _PR_ATTACH_MIME_TAG)
    if mime:
        return mime
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or "application/octet-stream"


# ---------------------------------------------------------------------------
# Local disk saving
# ---------------------------------------------------------------------------

# Extension → subfolder name under the Attachments root.
# Only extensions listed here are saved to disk during import.
_DISK_EXT_FOLDER = {
    ".pdf":  "pdf",
    ".doc":  "Word",  ".docx": "Word",
    ".xls":  "Excel", ".xlsx": "Excel", ".csv": "Excel",
    ".ppt":  "PowerPoint", ".pptx": "PowerPoint",
    ".mp4":  "Videos", ".avi":  "Videos", ".mov":  "Videos",
    ".wmv":  "Videos", ".mkv":  "Videos", ".flv":  "Videos",
    ".webm": "Videos", ".m4v":  "Videos", ".mpeg": "Videos",
    ".mpg":  "Videos", ".3gp":  "Videos",
    ".jpg":  "Images", ".jpeg": "Images", ".png":  "Images",
    ".gif":  "Images", ".bmp":  "Images", ".tiff": "Images",
    ".tif":  "Images", ".svg":  "Images", ".webp": "Images",
    ".txt":  "Text",   ".log":  "Text",   ".xml":  "Text",
    ".json": "Text",   ".html": "Text",   ".htm":  "Text",
}

def _save_to_disk(raw_data: bytes, filename: str, attach_dir: str,
                  email_date=None) -> str | None:
    """
    Save raw_data to attach_dir/<subfolder>/filename.
    Returns the path it was saved to, or None if the extension is not in the
    target list or saving fails.
    Duplicate filenames get a numeric suffix (_1, _2, …) rather than overwriting.
    If email_date is provided, the file's modification time is set to match.
    """
    ext = os.path.splitext(filename or "")[1].lower()
    subfolder = _DISK_EXT_FOLDER.get(ext)
    if not subfolder:
        return None                   # skip types not in the target list
    try:
        dest_dir = os.path.join(attach_dir, subfolder)
        os.makedirs(dest_dir, exist_ok=True)
        dest = os.path.join(dest_dir, filename)
        if os.path.exists(dest):
            base, fext = os.path.splitext(filename)
            n = 1
            while os.path.exists(dest):
                dest = os.path.join(dest_dir, f"{base}_{n}{fext}")
                n += 1
        with open(dest, "wb") as fh:
            fh.write(raw_data)
        # Set file timestamps to match the email date
        if email_date is not None:
            try:
                import calendar
                if hasattr(email_date, 'timestamp'):
                    ts = email_date.timestamp()
                else:
                    ts = calendar.timegm(email_date.timetuple())
                os.utime(dest, (ts, ts))
            except Exception:
                pass  # non-fatal — file is saved, timestamp just won't match
        return dest
    except Exception as e:
        print(f"  WARNING: could not save {filename} to disk: {e}", flush=True)
        return None


def extract_attachments(message, email_id, fs=None, attach_dir=None, email_date=None):
    """
    Return attachment metadata list.
    If fs (GridFS instance) is provided, binary data is stored there
    and each entry gets a 'gridfs_id' field for later download.
    If attach_dir is provided, pdf/Word/Excel files are also saved to disk.
    """
    attachments = []
    try:
        count = message.number_of_attachments
    except Exception:
        return attachments

    for i in range(count):
        meta = {"index": i, "size_bytes": 0, "filename": f"attachment_{i}", "gridfs_id": None}
        try:
            att          = message.get_attachment(i)
            size         = att.size or 0
            filename     = _att_filename(att, i)
            content_type = _att_mime(att, filename)
            meta.update({"size_bytes": size, "filename": filename, "content_type": content_type})

            if size > 0:
                source_key = f"{email_id}::{i}"
                # Read binary once — used for GridFS storage and text extraction
                try:
                    raw_data = att.read_buffer(size)
                except Exception as e:
                    raw_data = None
                    meta["read_error"] = str(e)

                # Store in GridFS (with dedup)
                if fs is not None and raw_data is not None:
                    existing = fs._files.find_one(
                        {"metadata.source_key": source_key}, {"_id": 1}
                    )
                    if existing:
                        meta["gridfs_id"] = str(existing["_id"])
                    else:
                        try:
                            file_id = fs.put(
                                raw_data,
                                filename=filename,
                                content_type=content_type,
                                metadata={"source_key": source_key, "email_id": email_id},
                            )
                            meta["gridfs_id"] = str(file_id)
                        except Exception as e:
                            meta["gridfs_error"] = str(e)

                # Save to local disk (pdf / Word / Excel only)
                if attach_dir and raw_data is not None:
                    saved_path = _save_to_disk(raw_data, filename, attach_dir,
                                               email_date=email_date)
                    if saved_path:
                        meta["disk_path"] = saved_path

                # Extract searchable text from attachment
                if raw_data is not None:
                    extracted = extract_text_from_bytes(raw_data, filename, content_type)
                    if extracted.strip():
                        meta["extracted_text"] = extracted
        except Exception as e:
            meta["error"] = str(e)

        attachments.append(meta)
    return attachments


def get_item_type(message, folder_path):
    """
    Determine item type: email | contact | appointment | task | other.
    Tries the pypff message_class property first, falls back to folder name.
    """
    try:
        mc = safe_str(getattr(message, "message_class", "")).strip().upper()
        if mc.startswith("IPM.CONTACT"):
            return "contact"
        if mc.startswith("IPM.APPOINTMENT") or mc.startswith("IPM.SCHEDULE"):
            return "appointment"
        if mc.startswith("IPM.TASK"):
            return "task"
        if mc.startswith("IPM.NOTE") or mc == "IPM.":
            return "email"
        if mc:                         # any other explicit class
            return "other"
    except Exception:
        pass

    # Fall back to folder name keywords
    fp = folder_path.upper()
    if "CONTACT" in fp:
        return "contact"
    if "CALENDAR" in fp or "APPOINTMENT" in fp:
        return "appointment"
    if "TASK" in fp:
        return "task"
    return "email"


def message_to_doc(message, folder_path, index, fs=None, attach_dir=None):
    """Convert a pypff message into a MongoDB document."""
    subject = safe_str(message.subject)
    sender_name = safe_str(message.sender_name)

    # Body content
    body_plain = safe_str(message.plain_text_body)
    body_html  = safe_str(message.html_body)
    body_rtf   = safe_str(message.rtf_body)

    # Dates from pypff native properties
    delivery_time     = parse_date(message.delivery_time)
    creation_time     = parse_date(message.creation_time)
    modification_time = parse_date(message.modification_time)
    client_submit     = parse_date(message.client_submit_time)

    # Richer metadata from transport headers
    headers = parse_transport_headers(message)
    # Prefer header date if pypff delivery_time is missing
    best_date = delivery_time or client_submit or headers["date"]

    email_id    = make_id(folder_path, index)
    attachments = extract_attachments(message, email_id, fs, attach_dir=attach_dir,
                                      email_date=best_date)

    # Concatenate all attachment extracted texts for MongoDB text indexing
    attachment_text = "\n\n".join(
        a["extracted_text"] for a in attachments if a.get("extracted_text")
    )[:50_000]

    doc = {
        "_id": email_id,
        "folder_path":        folder_path,
        "item_index":         index,
        "subject":            subject,
        "sender_name":        sender_name,
        "from_addr":          headers["from_addr"],
        "to_addrs":           headers["to_addrs"],
        "cc_addrs":           headers["cc_addrs"],
        "message_id":         headers["message_id"],
        "conversation_topic": safe_str(message.conversation_topic),
        "date":               best_date,
        "delivery_time":      delivery_time,
        "creation_time":      creation_time,
        "modification_time":  modification_time,
        "client_submit_time": client_submit,
        # Truncate bodies to 50k chars to keep docs manageable
        "body_plain": body_plain[:50000] if body_plain else "",
        "body_html":  body_html[:50000]  if body_html  else "",
        "body_rtf":   body_rtf[:50000]   if body_rtf   else "",
        "has_attachments":  len(attachments) > 0,
        "attachments":      attachments,
        "attachment_text":  attachment_text,
        "item_type":        get_item_type(message, folder_path),
        "imported_at":      datetime.datetime.utcnow(),
    }
    return doc


# ---------------------------------------------------------------------------
# Folder walker
# ---------------------------------------------------------------------------

def _flush_batch(collection, docs, stats, batch):
    """Flush docs from the front of the list in place to bound memory."""
    while len(docs) >= batch:
        chunk = docs[:batch]
        del docs[:batch]
        stats["total"]       += len(chunk)
        stats["attachments"] += sum(1 for d in chunk for a in d.get("attachments", []) if a.get("gridfs_id"))
        stats["disk_saved"]  += sum(1 for d in chunk for a in d.get("attachments", []) if a.get("disk_path"))
        flush_to_mongo(collection, chunk, stats)
        print(f"  Flushed {stats['total']:>6} docs | new={stats['inserted']} updated={stats['updated']} att={stats['attachments']} disk={stats['disk_saved']} err={stats['errors']}", flush=True)


def walk_folder(folder, path, docs, verbose, fs=None, attach_dir=None, collection=None, stats=None, batch=100):
    folder_name  = safe_str(folder.name, fallback="(unnamed)")
    current_path = (path + "/" + folder_name) if path else folder_name

    # Messages in this folder
    msg_count = folder.number_of_sub_messages
    for i in range(msg_count):
        try:
            message = folder.get_sub_message(i)
            doc = message_to_doc(message, current_path, i, fs=fs, attach_dir=attach_dir)
            docs.append(doc)
            if verbose:
                ts = doc.get("date") or ""
                n_att = len([a for a in doc.get("attachments", []) if a.get("gridfs_id")])
                att_note = f" [{n_att} att]" if n_att else ""
                print(f"  [{current_path}] #{i}  {doc['subject'][:60]}  ({ts}){att_note}",
                      flush=True)
        except Exception as e:
            print(f"  WARNING: message {i} in '{current_path}': {e}", flush=True)
        if collection is not None and stats is not None and len(docs) >= batch:
            _flush_batch(collection, docs, stats, batch)

    # Recurse into sub-folders
    for i in range(folder.number_of_sub_folders):
        try:
            sub = folder.get_sub_folder(i)
            walk_folder(sub, current_path, docs, verbose, fs=fs, attach_dir=attach_dir, collection=collection, stats=stats, batch=batch)
        except Exception as e:
            print(f"  WARNING: sub-folder {i} of '{current_path}': {e}", flush=True)


# ---------------------------------------------------------------------------
# Bulk upsert
# ---------------------------------------------------------------------------

def flush_to_mongo(collection, docs, stats):
    if not docs:
        return
    ops = [UpdateOne({"_id": d["_id"]}, {"$set": d}, upsert=True) for d in docs]
    try:
        result = collection.bulk_write(ops, ordered=False)
        stats["inserted"] += result.upserted_count
        stats["updated"]  += result.modified_count
    except BulkWriteError as bwe:
        stats["errors"]   += len(bwe.details.get("writeErrors", []))
        stats["inserted"] += bwe.details.get("nUpserted", 0)
        stats["updated"]  += bwe.details.get("nModified", 0)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Search
# ---------------------------------------------------------------------------

def search(collection, query, limit):
    """Full-text search across subject and body_plain."""
    # Ensure text index exists (no-op if already created)
    collection.create_index(
        [("subject", "text"), ("body_plain", "text"), ("attachment_text", "text")],
        name="text_search",
        default_language="english",
        weights={"subject": 10, "body_plain": 5, "attachment_text": 3},
    )

    cursor = collection.find(
        {"$text": {"$search": query}},
        {"score": {"$meta": "textScore"}, "subject": 1, "from_addr": 1,
         "date": 1, "folder_path": 1, "body_plain": 1},
    ).sort([("score", {"$meta": "textScore"})]).limit(limit)

    results = list(cursor)
    if not results:
        print("No results found.")
        return

    print(f"{len(results)} result(s) for: {query!r}\n")
    print("-" * 70)
    for doc in results:
        print(f"Subject : {doc.get('subject', '(no subject)')}")
        print(f"From    : {doc.get('from_addr', '')}")
        print(f"Date    : {doc.get('date', '')}")
        print(f"Folder  : {doc.get('folder_path', '')}")
        body = (doc.get("body_plain") or "").strip()
        if body:
            # Show first 300 chars as a snippet
            snippet = " ".join(body.split())[:300]
            print(f"Snippet : {snippet}...")
        print("-" * 70)


def main():
    ap = argparse.ArgumentParser(
        description="Import PST/OST into MongoDB, or search imported emails."
    )
    ap.add_argument("--pst",        default=None,
                    help="Path to the .pst (or .ost) file (required for import)")
    ap.add_argument("--mongo",      default="mongodb://localhost:27017",
                    help="MongoDB connection string")
    ap.add_argument("--db",         default="pst_import",
                    help="MongoDB database name (default: pst_import)")
    ap.add_argument("--collection", default="pst_items",
                    help="MongoDB collection name (default: pst_items)")
    ap.add_argument("--batch",      type=int, default=100,
                    help="Bulk insert batch size (default: 100)")
    ap.add_argument("--verbose",    action="store_true",
                    help="Print each item as it is processed")
    ap.add_argument("--search",     default=None, metavar="QUERY",
                    help="Search subject and body; skips import")
    ap.add_argument("--limit",      type=int, default=20,
                    help="Max results to return for --search (default: 20)")
    ap.add_argument("--attach-dir", default=None, metavar="DIR",
                    help="Root folder for saving pdf/Word/Excel attachments to disk")
    args = ap.parse_args()

    # Connect to MongoDB
    try:
        client = MongoClient(args.mongo, serverSelectionTimeoutMS=5000)
        client.server_info()
        collection = client[args.db][args.collection]
    except Exception as e:
        print(f"ERROR: Cannot connect to MongoDB: {e}")
        sys.exit(1)

    # Search mode
    if args.search:
        search(collection, args.search, args.limit)
        return

    # Import mode — --pst is required
    if not args.pst:
        ap.error("--pst is required when not using --search")

    pst_path = os.path.abspath(args.pst)
    if not os.path.isfile(pst_path):
        print(f"ERROR: File not found: {pst_path}")
        sys.exit(1)

    print(f"PST file  : {pst_path}")
    print(f"MongoDB   : {args.mongo}")
    print(f"Database  : {args.db}.{args.collection}")
    print(f"Batch size: {args.batch}")
    print()
    print("MongoDB connection OK.")

    # Indexes for common query patterns
    collection.create_index("folder_path")
    collection.create_index("from_addr")
    collection.create_index("date")
    collection.create_index("subject")
    collection.create_index("message_id")
    collection.create_index("item_type")

    # Full-text search index — subject (weight 10), body (5), attachments (3)
    try:
        collection.create_index(
            [("subject", "text"), ("body_plain", "text"), ("attachment_text", "text")],
            name="text_search",
            default_language="english",
            weights={"subject": 10, "body_plain": 5, "attachment_text": 3},
        )
    except Exception:
        pass  # already exists; will be recreated if needed by the web app

    # GridFS instance for attachment storage
    fs = gridfs_module.GridFS(client[args.db])
    print("GridFS ready for attachments.")

    # Local disk folder for pdf / Word / Excel attachments
    attach_dir = os.path.abspath(args.attach_dir) if args.attach_dir else None
    if attach_dir:
        for sub in ("pdf", "Word", "Excel", "PowerPoint", "Videos", "Images", "Text"):
            os.makedirs(os.path.join(attach_dir, sub), exist_ok=True)
        print(f"Saving pdf/Word/Excel/Videos/Images/Text attachments to: {attach_dir}")

    # Open PST
    try:
        pst_file = pypff.file()
        pst_file.open(pst_path)
    except Exception as e:
        print(f"ERROR: Cannot open PST file: {e}")
        traceback.print_exc()
        sys.exit(1)

    root     = pst_file.get_root_folder()
    all_docs = []
    stats    = {"inserted": 0, "updated": 0, "errors": 0, "total": 0, "attachments": 0, "disk_saved": 0}

    print("Scanning folders...\n")
    for i in range(root.number_of_sub_folders):
        try:
            top = root.get_sub_folder(i)
            walk_folder(top, "", all_docs, args.verbose, fs=fs, attach_dir=attach_dir, collection=collection, stats=stats, batch=args.batch)
        except Exception as e:
            print(f"WARNING: top-level folder {i}: {e}", flush=True)

        # Flush whenever buffer exceeds batch size
        while len(all_docs) >= args.batch:
            batch    = all_docs[:args.batch]
            all_docs = all_docs[args.batch:]
            stats["total"]       += len(batch)
            stats["attachments"] += sum(
                1 for d in batch for a in d.get("attachments", []) if a.get("gridfs_id")
            )
            stats["disk_saved"] += sum(
                1 for d in batch for a in d.get("attachments", []) if a.get("disk_path")
            )
            flush_to_mongo(collection, batch, stats)
            print(f"  Flushed {stats['total']:>6} docs | "
                  f"new={stats['inserted']} updated={stats['updated']} "
                  f"att={stats['attachments']} disk={stats['disk_saved']} err={stats['errors']}",
                  flush=True)

    # Final flush
    if all_docs:
        stats["total"]       += len(all_docs)
        stats["attachments"] += sum(
            1 for d in all_docs for a in d.get("attachments", []) if a.get("gridfs_id")
        )
        stats["disk_saved"] += sum(
            1 for d in all_docs for a in d.get("attachments", []) if a.get("disk_path")
        )
        flush_to_mongo(collection, all_docs, stats)

    pst_file.close()

    print()
    print("=" * 50)
    print("Import complete.")
    print(f"  Total processed : {stats['total']}")
    print(f"  Newly inserted  : {stats['inserted']}")
    print(f"  Updated (dupes) : {stats['updated']}")
    print(f"  Attachments     : {stats['attachments']}")
    print(f"  Errors          : {stats['errors']}")
    print(f"  Collection      : {args.db}.{args.collection}")
    print("=" * 50)


if __name__ == "__main__":
    main()